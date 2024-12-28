import io
import json
import os
import tempfile
import uuid
import time
import base64
import zlib
from urllib.parse import unquote
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
import requests
import boto3
import sys
import base64
from botocore.exceptions import ClientError
from utils import refresh_user_google, refresh_user_dropbox, lambda_timelimit_exceeded, lambda_time_left_seconds, get_user_table_entry, extend_lock_time, set_user_table_cache_entry, prtime
from typing import Dict, List, Tuple, Any, Optional, Union
from dataclasses import dataclass
import datetime
from datetime import timezone
import subprocess
import dropbox
import jsons
import gzip
from pdf import read_pdf
import openpyxl
from bs4 import BeautifulSoup
import re
import zipfile
import tarfile
import traceback
from lock import lock_user_local, unlock_user_local
import shutil
import mimetypes
from typing import TYPE_CHECKING
if TYPE_CHECKING:
    from mypy_boto3_s3.client import S3Client
else:
    S3Client = object

from googleapiclient.http import MediaIoBaseDownload
from google.api_core.datetime_helpers import to_rfc3339, from_rfc3339, from_microseconds
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
import google.oauth2.credentials 
import google.auth.transport.requests
import googleapiclient.discovery

import docx
from pptx import Presentation

from distilbert_dotprod import MsmarcoDistilbertBaseDotProdV3
import pickle
from faiss_rm import FaissRM, DocStorageType
from custom_model import CustomModel
from text_utils import format_paragraph

vectorizer_cache = {}
def get_vectorizer(email, tracebuf):
    global vectorizer_cache
    if email in vectorizer_cache:
        print(f"get_vectorizer: Returning {type(vectorizer_cache[email])} object from cache for {email}")
        return vectorizer_cache[email]
    item = get_user_table_entry(email)
    if 'CustomModelBucket' in item and 'CustomModelObjectKey' in item:
        vectorizer_cache[email] = CustomModel(item['CustomModelBucket']['S'], item['CustomModelObjectKey']['S'], tracebuf)
        print(f"get_vectorizer: CustomModelBucket {item['CustomModelBucket']['S']} CustomModelObjectKey {item['CustomModelObjectKey']['S']}. Returning {type(vectorizer_cache[email])} object from cache for {email}")
    elif os.path.isdir('/var/task/sentence-transformers/msmarco-distilbert-base-dot-prod-v3'):
        vectorizer_cache[email] = MsmarcoDistilbertBaseDotProdV3(
                tokenizer_name_or_path='/var/task/sentence-transformers/msmarco-distilbert-base-dot-prod-v3',
                model_name_or_path='/var/task/sentence-transformers/msmarco-distilbert-base-dot-prod-v3'
            )
        print(f"get_vectorizer: isdir /var/task/sentence-transformers/msmarco-distilbert-base-dot-prod-v3. Returning {type(vectorizer_cache[email])} object from cache for {email}")
    else:
        vectorizer_cache[email] = MsmarcoDistilbertBaseDotProdV3()
        print(f"get_vectorizer: No custom model and no model in /var. Returning {type(vectorizer_cache[email])} object from cache for {email}")
    return vectorizer_cache[email]

def lock_user(email, index_dir, client, takeover_lock_end_time=0):
    print(f"lock_user_aws: Entered. Trying to lock for email={email}")
    item = get_user_table_entry(email)
    print(f"lock_user_aws: User table entry for {email}={item}")
    now = time.time()
    now_s = datetime.datetime.fromtimestamp(now).strftime('%Y-%m-%d %I:%M:%S')
    gdrive_next_page_token = None
    dropbox_next_page_token = None
    if not 'lock_end_time' in item:
        print(f"lock_user_aws: no lock_end_time in ddb. now={now}/{now_s}")
        l_e_t = None
        l_e_t_s = None
    else:
        l_e_t=int(float(item['lock_end_time']['N']))
        l_e_t_s = datetime.datetime.fromtimestamp(l_e_t).strftime('%Y-%m-%d %I:%M:%S')
        print(f"lock_user_aws: lock_end_time in ddb={l_e_t}/{l_e_t_s}, now={now}/{now_s}")
        if takeover_lock_end_time != 0:
            if takeover_lock_end_time == l_e_t:
                if 'gdrive_next_page_token' in item:
                    gdrive_next_page_token = item['gdrive_next_page_token']['S']
                if 'dropbox_next_page_token' in item:
                    dropbox_next_page_token = item['dropbox_next_page_token']['S']
                print(f"lock_user_aws: Takeover successful for {email}. lock_end_time {l_e_t}/{l_e_t_s}")
                return gdrive_next_page_token, dropbox_next_page_token, True
            else:
                print(f"lock_user_aws: Takeover unsuccessful for {email}. ddb lock_end_time {l_e_t}/{l_e_t_s}, takeover_lock_end_time={takeover_lock_end_time}")
                return gdrive_next_page_token, dropbox_next_page_token, False
        else:
            print(f"lock_user_aws: takeover_lock_end_time is 0. Not attempting to take over existing lock. Proceeding with new lock attempt")
    time_left = int(lambda_time_left_seconds())
    if time_left > (60 * 15):
        time_left = 60*15
    if index_dir:
        if lock_user_local(index_dir):
            ute = {'email': {'S': email}, 'lock_end_time': {'N': str(int(now)+(time_left))}}
            set_user_table_cache_entry(email, ute)
            return 'notused_gdrive_next_token', 'notused_dropbox_next_token', True
        else:
            return 'notused_gdrive_next_token', 'notused_dropbox_next_token', False
    else:
        try:
            response = client.update_item(
                TableName=os.environ['USERS_TABLE'],
                Key={'email': {'S': email}},
                UpdateExpression="set #lm = :st",
                ConditionExpression=f"attribute_not_exists(#lm) OR #lm < :nw",
                ExpressionAttributeNames={'#lm': 'lock_end_time'},
                ExpressionAttributeValues={':nw': {'N': str(int(now))}, ':st': {'N': str(int(now)+(time_left))} },
                ReturnValues="ALL_NEW"
            )
            if 'Attributes' in response:
                if 'gdrive_next_page_token' in response['Attributes']:
                    gdrive_next_page_token = response['Attributes']['gdrive_next_page_token']['S']
                if 'dropbox_next_page_token' in response['Attributes']:
                    dropbox_next_page_token = response['Attributes']['dropbox_next_page_token']['S']
                set_user_table_cache_entry(email, response['Attributes'])
            print(f"lock_user_aws: conditional check success. No other instance of lambda is active for user {email}. "
                f"gdrive_next_page_token={gdrive_next_page_token}, dropbox_next_page_token={dropbox_next_page_token}")
            return gdrive_next_page_token, dropbox_next_page_token, True
        except ClientError as e:
            if e.response['Error']['Code'] == "ConditionalCheckFailedException":
                print(f"lock_user_aws: conditional check failed. {e.response['Error']['Message']}. Another instance of lambda is active for {email}")
                return gdrive_next_page_token, dropbox_next_page_token, False
            else:
                raise

def update_next_page_tokens(email, client, gdrive_next_page_token, dropbox_next_page_token):
    print(f"update_next_page_tokens: Entered. email={email}")
    if 'USERS_TABLE' not in os.environ:
        print(f"update_next_page_tokens: USERS_TABLE environment variable not present. Assume local execution")
        return True
    try:
        if gdrive_next_page_token and dropbox_next_page_token:
            ue="SET gdrive_next_page_token = :pt, dropbox_next_page_token = :db"
            eav={':pt': {'S': gdrive_next_page_token}, ':db': {'S': dropbox_next_page_token}}
            response = client.update_item(
                TableName=os.environ['USERS_TABLE'],
                Key={'email': {'S': email}},
                UpdateExpression=ue, ExpressionAttributeValues=eav,
                ReturnValues="ALL_NEW"
            )
        elif gdrive_next_page_token and not dropbox_next_page_token:
            ue="SET gdrive_next_page_token = :pt  REMOVE dropbox_next_page_token"
            eav={':pt': {'S': gdrive_next_page_token}}
            response = client.update_item(
                TableName=os.environ['USERS_TABLE'],
                Key={'email': {'S': email}},
                UpdateExpression=ue, ExpressionAttributeValues=eav,
                ReturnValues="ALL_NEW"
            )
        elif not gdrive_next_page_token and dropbox_next_page_token:
            ue="REMOVE gdrive_next_page_token  SET dropbox_next_page_token = :db"
            eav={':db': {'S': dropbox_next_page_token}}
            response = client.update_item(
                TableName=os.environ['USERS_TABLE'],
                Key={'email': {'S': email}},
                UpdateExpression=ue, ExpressionAttributeValues=eav,
                ReturnValues="ALL_NEW"
            )
        else:
            ue = "REMOVE gdrive_next_page_token, dropbox_next_page_token"
            response = client.update_item(
                TableName=os.environ['USERS_TABLE'],
                Key={'email': {'S': email}},
                UpdateExpression=ue,
                ReturnValues="ALL_NEW"
            )
        set_user_table_cache_entry(email, response['Attributes'])
    except ClientError as e:
        print(f"update_next_page_tokens: Error {e.response['Error']['Message']} while unlocking")
        return False
    return True

def download_gdrive_file(service, file_id, filename) -> io.BytesIO:
  """Downloads a file
  Args:
      file_id: ID of the file to download
  Returns : IO object with location.

  Load pre-authorized user credentials from the environment.
  TODO(developer) - See https://developers.google.com/identity
  for guides on implementing OAuth2 for the application.
  """
  try:
    # pylint: disable=maybe-no-member
    request = service.files().get_media(fileId=file_id)
    file = io.BytesIO()
    downloader = MediaIoBaseDownload(file, request)
    done = False
    while done is False:
      status, done = downloader.next_chunk()
      print(f"filename={filename}, fileid={file_id}, download {int(status.progress() * 100)} %")
  except HttpError as error:
    print(f"download_gdrive_file: An error occurred: {error}")
    return None
  file.seek(0, os.SEEK_SET)
  return file

def export_gdrive_file(access_token, file_id, fmt) -> io.BytesIO:
    headers={"Authorization": f"Bearer {access_token}"}
    if fmt == 'pptx':
        url = f"https://docs.google.com/presentation/d/{file_id}/export/{fmt}"
    elif fmt == 'docx':
        url = f"https://docs.google.com/document/d/{file_id}/export?format=doc"
    elif fmt == 'xlsx':
        url = f"https://docs.google.com/spreadsheets/d/{file_id}/export?format=xlsx"
    else:
        print(f"export_gdrive_file: Unknown format {fmt} for fileid {file_id}. Do not know how to download..")
        return None
    for attempt in range(2):
        try:
            resp = requests.get(url, stream=True, headers=headers)
            resp.raise_for_status()
            file = io.BytesIO()
            for chunk in resp.iter_content(chunk_size=1024): 
                if chunk:
                    file.write(chunk)
            file.seek(0, os.SEEK_SET)
            # open the files. If the open throws an exception, sleep and retry one time
            if fmt == 'xlsx':
                wb = openpyxl.load_workbook(file)
                file.seek(0, os.SEEK_SET)
            elif fmt == 'pptx':
                prs = Presentation(file)
                file.seek(0, os.SEEK_SET)
            elif fmt == 'docx':
                doc = docx.Document(file)
                file.seek(0, os.SEEK_SET)
            return file
        except requests.exceptions.HTTPError as http_err:
            if resp.status_code == 401:
                print(f"401 Unauthorized error occurred for file_id {file_id}")
                return None
            else:
                print(f"HTTP error occurred: {http_err}. Sleeping and continuing..")
        except Exception as ex:
            print(f"export_gdrive_file: Caught {ex}. Sleeping 10 seconds")
            time.sleep(10)
    return None

@dataclass
class IndexMetadata:
    jsonl_last_modified:float  = 0 # time.time()
    index_flat_last_modified:float = 0
    index_ivfadc_last_modified: float = 0
    bm25s_index_last_modified: float = 0
    
    def is_vdb_index_stale(self):
        if not self.jsonl_last_modified and not self.index_flat_last_modified and \
                not self.index_ivfadc_last_modified and not self.bm25s_index_last_modified:
            print(f"IndexMetadata.is_vdb_index_stale: all values 0. Returning stale")
            return True
        return self.jsonl_last_modified > self.index_flat_last_modified \
                or self.jsonl_last_modified > self.index_ivfadc_last_modified \
                or self.jsonl_last_modified > self.bm25s_index_last_modified

FILES_INDEX_JSONL = "/tmp/files_index.jsonl"    
FILES_INDEX_JSONL_GZ = "/tmp/files_index.jsonl.gz"
INDEX_METADATA_JSON = "/tmp/index_metadata.json"
FAISS_INDEX_FLAT = "/tmp/faiss_index_flat"
FAISS_INDEX_IVFADC = "/tmp/faiss_index_ivfadc"
BM25S_INDEX = "/tmp/bm25s_index"
BM25S_INDEX_TAR = "/tmp/bm25s_index.tar"

def make_temp_copy(index_dir, copy_faiss_index):
    src = os.path.join(index_dir, os.path.basename(FILES_INDEX_JSONL_GZ))
    dst = FILES_INDEX_JSONL_GZ
    try:
        shutil.copy2(src, dst)
    except Exception as ex:
        print(f"make_temp_copy: Caught {ex} copying {src} to {dst}")
        return False
    src = os.path.join(index_dir, os.path.basename(INDEX_METADATA_JSON))
    dst = INDEX_METADATA_JSON
    try:
        shutil.copy2(src, dst)
    except Exception as ex:
        print(f"make_temp_copy: Caught {ex} copying {src} to {dst}")
        return False
    if copy_faiss_index:
        src = os.path.join(index_dir, os.path.basename(FAISS_INDEX_FLAT))
        dst = FAISS_INDEX_FLAT
        try:
            shutil.copy2(src, dst)
        except Exception as ex:
            print(f"make_temp_copy: Caught {ex} copying {src} to {dst}")
            return False
        src = os.path.join(index_dir, os.path.basename(FAISS_INDEX_IVFADC))
        dst = FAISS_INDEX_IVFADC
        try:
            shutil.copy2(src, dst)
        except Exception as ex:
            print(f"make_temp_copy: Caught {ex} copying {src} to {dst}")
            return False
        src = os.path.join(index_dir, os.path.basename(BM25S_INDEX_TAR))
        dst = BM25S_INDEX_TAR
        try:
            shutil.copy2(src, dst)
        except Exception as ex:
            print(f"make_temp_copy: Caught {ex} copying {src} to {dst}")
            return False
    return True

def download_files_index(s3client, bucket, prefix, download_faiss_index) -> bool:
    """ download the jsonl file that has a line for each file in the google drive; download it to /tmp/files_index.jsonl   
        download_faiss_index:  download the built faiss indexes to local storage (in addition to above jsonl file)
        return True if all files can be downloaded.  Return False if any of the files cannot be downloaded.  
        Note that files_index.jsonl.gz may exist but faiss_index_flat and faiss_index_ivfadc may not exist since the latter is sometimes built only after all the document embeddings are generated. """
    try:
        files_index_jsonl_gz = FILES_INDEX_JSONL_GZ
        # index1/raj@yoja.ai/files_index.jsonl.gz
        print(f"Downloading {files_index_jsonl_gz} from s3://{bucket}/{prefix}")
        s3client.download_file(bucket, f"{prefix}/{os.path.basename(files_index_jsonl_gz)}", files_index_jsonl_gz)
    except Exception as ex:
        print(f"Caught {ex} while downloading {files_index_jsonl_gz} from s3://{bucket}/{prefix}")
        return False

    try:
        index_metadata_fname = INDEX_METADATA_JSON
        print(f"Downloading {index_metadata_fname} from s3://{bucket}/{prefix}")
        s3client.download_file(bucket, f"{prefix}/{os.path.basename(index_metadata_fname)}", index_metadata_fname)
    except Exception as ex:
        print(f"Caught {ex} while downloading {index_metadata_fname}.  Creating an empty {index_metadata_fname}")
        # create an empty metadata file
        _create_index_metadata_json_local(IndexMetadata())

    if download_faiss_index:
        try:
            faiss_index_flat_fname = FAISS_INDEX_FLAT
            print(f"Downloading {faiss_index_flat_fname} from s3://{bucket}/{prefix}")
            # index1/raj@yoja.ai/faiss_index_flat
            s3client.download_file(bucket, f"{prefix}/{os.path.basename(faiss_index_flat_fname)}", faiss_index_flat_fname)
        except Exception as ex:
            print(f"Caught {ex} while downloading {faiss_index_flat_fname} from s3://{bucket}/{prefix}")
            return False
        try:
            faiss_index_ivfadc_fname = FAISS_INDEX_IVFADC
            print(f"Downloading {faiss_index_ivfadc_fname} from s3://{bucket}/{prefix}")
            # index1/raj@yoja.ai/faiss_index_ivfadc
            s3client.download_file(bucket, f"{prefix}/{os.path.basename(faiss_index_ivfadc_fname)}", faiss_index_ivfadc_fname)
        except Exception as ex:
            print(f"Caught {ex} while downloading {faiss_index_ivfadc_fname} from s3://{bucket}/{prefix}")
            return False
        try:
            print(f"Downloading {BM25S_INDEX_TAR} from s3://{bucket}/{prefix}")
            # index1/raj@yoja.ai/bm25s_index
            s3client.download_file(bucket, f"{prefix}/{os.path.basename(BM25S_INDEX_TAR)}", BM25S_INDEX_TAR)
        except Exception as ex:
            print(f"Caught {ex} while downloading {BM25S_INDEX_TAR} from s3://{bucket}/{prefix}")
            return False
    
    return True

def _write_progress_file(index_dir, s3client, bucket, prefix, attr_prefix, num_unmodified, unmodified_size, num_needs_embedding, needs_embedding_size):
    try:
        with open('/tmp/indexing_progress.json', 'w+') as fl:
            fl.write(json.dumps({f"{attr_prefix}_num_unmodified": num_unmodified,
                                 f"{attr_prefix}_unmodified_size": unmodified_size,
                                 f"{attr_prefix}_num_needs_embedding": num_needs_embedding,
                                 f"{attr_prefix}_needs_embedding_size": needs_embedding_size}))
        if index_dir:
            shutil.copy2('/tmp/indexing_progress.json', os.path.join(index_dir, 'indexing_progress.json'))
        else:
            s3client.upload_file('/tmp/indexing_progress.json', bucket, f"{prefix}/indexing_progress.json")
    except Exception as ex:
        print(f"Caught {ex} while writing progress file")
        return False

def init_vdb(email, index_dir, s3client, bucket, prefix, doc_storage_type:DocStorageType,
            chat_config=None, tracebuf=None, build_faiss_indexes=True, sub_prefix=None) -> FaissRM :
    """
    initializes a faiss vector db with the embeddings specified in bucket/prefix/files_index.jsonl.
    Downloads the index from S3.  Returns a FaissRM instance which encapsulates vectorDB, metadata, documents.  Or None, if index not found in S3
    sub_prefix: specify subfolder under which index must be downloaded from.  If not specified, ignored.
    """
    print(f"init_vdb: Entered. email={email}, index_dir={index_dir}, bucket={bucket}, prefix={prefix}, sub_prefix={sub_prefix}")
    user_prefix = f"{prefix}/{email}" + f"{'/' + sub_prefix if sub_prefix else ''}"
    tracebuf.append(f"{prtime()} init_vdb: Entered. index_dir={index_dir}, bucket={bucket}, prefix={user_prefix}")
    fls = {}
    embeddings = []
    index_map = [] # list of (fileid, paragraph_index)
    flat_index_fname=None if build_faiss_indexes else FAISS_INDEX_FLAT
    ivfadc_index_fname=None if build_faiss_indexes else FAISS_INDEX_IVFADC
    bm25s_index_fname=None if build_faiss_indexes else BM25S_INDEX_TAR
    if index_dir:
        sts = make_temp_copy(index_dir, not build_faiss_indexes)
    else:
        sts = download_files_index(s3client, bucket, user_prefix, not build_faiss_indexes)
    if sts:
        with gzip.open(FILES_INDEX_JSONL_GZ, "r") as rfp:
            for line in rfp:
                finfo = json.loads(line)
                finfo['mtime'] = from_rfc3339(finfo['mtime'])
                fls[finfo['fileid']] = finfo
                if 'slides' in finfo:
                    key = 'slides'
                elif 'paragraphs' in finfo:
                    key = 'paragraphs'
                elif 'rows' in finfo:
                    key = 'rows'
                else:
                    continue
                for para_index in range(len(finfo[key])):
                    para = finfo[key][para_index]
                    if para:
                        if not flat_index_fname:
                            embeddings.append(pickle.loads(base64.b64decode(para['embedding'].strip()))[0])
                        del para['embedding']
                        index_map.append((finfo['fileid'], para_index))
    else:
        print(f"init_vdb: Failed to download files_index.jsonl from s3://{bucket}/{user_prefix}")
        return None
    print(f"init_vdb: finished reading files_index.jsonl.gz. Num files={len(fls.items())}")
    print(f"init_vdb: finished loading index_map. Entries in index_map={len(index_map)}")
    print(f"init_vdb: finished loading embeddings. Entries in embeddings={len(embeddings)}")
    tracebuf.append(f"{prtime()} init_vdb: finished loading files_index")
    vectorizer = get_vectorizer(email, tracebuf)
    return FaissRM(fls, index_map, embeddings, vectorizer, doc_storage_type,
                    chat_config, tracebuf, k=100,
                    flat_index_fname=flat_index_fname, ivfadc_index_fname=ivfadc_index_fname,
                    bm25s_index_fname=bm25s_index_fname)

def _calc_path(service, entry, folder_details):
    # calculate path by walking up parents
    path = ""
    while True:
        if not 'parents' in entry:
            break
        folder_id = entry['parents'][0]
        if folder_id not in folder_details:
            try:
                res = service.files().get(fileId=folder_id, fields='name,parents').execute()
                if 'parents' in res:
                    folder_details[folder_id] = {'filename': res['name'], 'parents': res['parents']}
                else:
                    folder_details[folder_id] = {'filename': res['name']}
            except Exception as ex:
                print(f"_calc_path: Exception {ex} while getting shared drive name for id {folder_id}")
        if folder_id in folder_details:
            dir_entry = folder_details[folder_id]
            path = dir_entry['filename'] + '/' + path
            if 'parents' in dir_entry:
                entry = dir_entry
                continue
            else:
                break
        else:
            print(f"_calc_path: WARNING: Could not map folder id {folder_id} to folder_name for filename {entry['filename']}, path={path}")
            break
    return path

def calc_file_lists(service, s3_index, gdrive_listing, folder_details) -> Tuple[dict, dict, dict]:
    """ returns a tuple of (unmodified:dict, needs_embedding:dict, deleted:dict).  Each dict has the format { fileid: {filename:abc, fileid:xyz, mtime:mno}}"""
    unmodified = {}
    deleted = {}
    needs_embedding = {}
    for fileid, gdrive_entry in gdrive_listing.items():
        if fileid in s3_index:
            if s3_index[fileid]['mtime'] == gdrive_entry['mtime']:
                if 'partial' in s3_index[fileid]:
                    #print(f"calc_file_lists: entry {gdrive_entry['filename']} is in s3 listing. Needs embedding. Time matches, but PARTIAL. gdrive_entry={gdrive_entry['mtime']}, s3listing={s3_index[fileid]['mtime']}")
                    needs_embedding[fileid] = s3_index[fileid]
                    needs_embedding[fileid]['mtime'] = gdrive_entry['mtime']
                else:
                    #print(f"calc_file_lists: entry {gdrive_entry['filename']} is in s3 listing. Unmodified. Time matches, no PARTIAL. gdrive_entry={gdrive_entry['mtime']}, s3listing={s3_index[fileid]['mtime']}")
                    unmodified[fileid] = s3_index[fileid]
            else:
                #print(f"calc_file_lists: entry {gdrive_entry['filename']} is in s3 listing. Needs embedding. Time mismatch. gdrive_entry={gdrive_entry['mtime']}, s3listing={s3_index[fileid]['mtime']}")
                needs_embedding[fileid] = s3_index[fileid]
                needs_embedding[fileid]['mtime'] = gdrive_entry['mtime']
        else:
            #print(f"calc_file_lists: entry {gdrive_entry['filename']} is not in s3 listing. Needs embedding")
            path=_calc_path(service, gdrive_entry, folder_details)
            gdrive_entry['path'] = path
            needs_embedding[fileid] = gdrive_entry
    
    # detect deleted files
    for fileid in s3_index:
        if not ( unmodified.get(fileid) or needs_embedding.get(fileid) ): 
            deleted[fileid] = s3_index[fileid]
            
    return unmodified, needs_embedding, deleted

def get_s3_index(index_dir, s3client, bucket, prefix) -> Dict[str, dict]:
    rv = {}
    if index_dir:
        sts = make_temp_copy(index_dir, False)
    else:
        sts = download_files_index(s3client, bucket, prefix, False)
    if sts:
        with gzip.open(FILES_INDEX_JSONL_GZ, "rb") as rfp:
            for line in rfp:
                ff = json.loads(line)
                ff['mtime'] = from_rfc3339(ff['mtime'])
                rv[ff['fileid']] = ff
    else:
        print(f"get_s3_index: Failed to download files_index.jsonl from s3://{bucket}/{prefix}")
    return rv

def read_docx(email, filename:str, fileid:str, file_io:io.BytesIO, mtime:datetime.datetime, prev_paras) -> Dict[str, Union[str,Dict[str, str]]]:
    doc_dct={"filename": filename, "fileid": fileid, "mtime": mtime, "paragraphs": prev_paras}
    prev_len = len(prev_paras)
    doc = docx.Document(file_io)
    vectorizer = get_vectorizer(email, [])
    prelude = f"The filename is {filename} and the paragraphs are:"
    prelude_token_len = vectorizer.get_token_count(prelude)
    chunk_len = prelude_token_len
    chunk_paras = []
    for parag in doc.paragraphs:
        para = parag.text
        if para:
            para_len = vectorizer.get_token_count(para)
            if para_len + chunk_len >= 512:
                if prev_len > len(doc_dct['paragraphs']): # skip previously processed chunks
                    chunk_len = prelude_token_len
                    chunk_paras = []
                    continue
                chunk = '.'.join(chunk_paras)
                para_dct = {'paragraph_text': chunk} # {'paragraph_text': 'Module 2: How To Manage Change', 'embedding': 'gASVCBsAAAAAAA...GVhLg=='}
                try:
                    embedding = vectorizer([f"{prelude}{chunk}"])
                    eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                    para_dct['embedding'] = eem
                except Exception as ex:
                    print(f"Exception {ex} while creating para embedding")
                doc_dct['paragraphs'].append(para_dct)
                chunk_len = prelude_token_len
                chunk_paras = []
                if lambda_timelimit_exceeded():
                    doc_dct['partial'] = "true"
                    print(f"read_docx: Lambda timelimit exceeded when reading docx file. Breaking..")
                    break
            else:
                chunk_paras.append(para)
                chunk_len += para_len
    if chunk_len > prelude_token_len and len(chunk_paras) > 0:
        if prev_len <= len(doc_dct['paragraphs']): # skip previously processed chunks
            chunk = '.'.join(chunk_paras)
            para_dct = {} # {'paragraph_text': 'Module 2: How To Manage Change', 'embedding': 'gASVCBsAAAAAAA...GVhLg=='}
            para_dct['paragraph_text'] = chunk
            try:
                embedding = vectorizer([f"{prelude}{chunk}"])
                eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                para_dct['embedding'] = eem
            except Exception as ex:
                print(f"Exception {ex} while creating residual para embedding")
            doc_dct['paragraphs'].append(para_dct)
    return doc_dct

def read_xlsx(email, filename, fileid, file_io, mtime:datetime.datetime, prev_rows) -> Dict[str, Union[str, Dict[str,str]]]:
    vectorizer = get_vectorizer(email, [])
    workbook = {'filename': filename, 'rows': prev_rows}
    wb = openpyxl.load_workbook(file_io)
    for sn in wb.sheetnames:
        sheet=wb[sn]
        row_in_files_index_jsonl = 0
        for row in range(1, sheet.max_row+1):
            rowstrings = ""
            for cl in range(1, sheet.max_column+1):
                cell = sheet.cell(row, cl)
                if cell.value and isinstance(cell.value, str):
                    rowstrings = rowstrings + cell.value.strip().rstrip('.') + '.'
            if rowstrings:
                if rowstrings[0] != '=':
                    row_in_files_index_jsonl += 1
                    if row_in_files_index_jsonl >= len(prev_rows):
                        row_dct = {'text': rowstrings, 'row_number': row}
                        chu = f"The file is named '{filename}', the sheet in the file is named '{sn}' and the row contains '{rowstrings}'"
                        try:
                            row_dct['embedding'] = base64.b64encode(pickle.dumps(vectorizer([chu]))).decode('ascii')
                        except Exception as ex:
                            print(f"Exception {ex} while creating embedding for row")
                        workbook['rows'].append(row_dct)
            if lambda_timelimit_exceeded():
                workbook['partial'] = "true"
                print(f"read_xlsx: Lambda timelimit exceeded reading xlsx file. Breaking..")
                break
    return workbook

def read_pptx(email, filename, fileid, file_io, mtime:datetime.datetime, prev_slides) -> Dict[str, Union[str, Dict[str,str]]]:
    vectorizer = get_vectorizer(email, [])
    prs = Presentation(file_io)
    ppt={"filename": filename, "fileid": fileid, "mtime": mtime, "slides": prev_slides}
    ind = 0
    for slide in prs.slides:
        title=None
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') \
                                and hasattr(slide.shapes.title, 'text'):
            title=slide.shapes.title.text
        txt=[]
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text != title:
                        txt.append(run.text)
        slide_text = ','.join(txt)
        if title:
            slide_dct = {"title": title, "text": slide_text}
            chu = f"The title of this slide is {title} and the content of the slide is {slide_text}"
        else:
            slide_dct = {"text": slide_text}
            chu = f"The content of the slide is {slide_text}"
        if ind >= len(prev_slides): # skip past prev_slides
            embedding = vectorizer([chu])
            try:
                eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                slide_dct['embedding'] = eem
            except Exception as ex:
                print(f"Exception {ex} while creating slide embedding")
            ppt['slides'].append(slide_dct)
            if lambda_timelimit_exceeded():
                ppt['partial'] = "true"
                print(f"read_pptx: Lambda timelimit exceeded reading pptx file. Breaking..")
                break
        ind += 1
    return ppt

def read_txt(email, filename:str, fileid:str, file_io:io.BytesIO, mtime:datetime.datetime, prev_paras) -> Dict[str, Union[str,Dict[str, str]]]:
    vectorizer = get_vectorizer(email, [])
    doc_dct={"filename": filename, "fileid": fileid, "mtime": mtime, "paragraphs": prev_paras}
    prev_len = len(prev_paras)
    fulltxt = file_io.getvalue().decode('utf-8')
    sentences = fulltxt.split('.')
    prelude = f"The filename is {filename} and the paragraphs are:"
    prelude_token_len = vectorizer.get_token_count(prelude)
    chunk_len = prelude_token_len
    chunk_paras = []
    for para in sentences:
        if para:
            para_len = vectorizer.get_token_count(para)
            if para_len + chunk_len >= 512:
                if prev_len > len(doc_dct['paragraphs']): # skip previously processed chunks
                    chunk_len = prelude_token_len
                    chunk_paras = []
                    continue
                chunk = '.'.join(chunk_paras)
                para_dct = {'paragraph_text': chunk}
                try:
                    embedding = vectorizer([f"{prelude}{chunk}"])
                    eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                    para_dct['embedding'] = eem
                except Exception as ex:
                    print(f"Exception {ex} while creating para embedding")
                doc_dct['paragraphs'].append(para_dct)
                chunk_len = prelude_token_len
                chunk_paras = []
                if lambda_timelimit_exceeded():
                    doc_dct['partial'] = "true"
                    print(f"read_txt: Lambda timelimit exceeded when reading docx file. Breaking..")
                    break
            else:
                chunk_paras.append(para)
                chunk_len += para_len
    if chunk_len > prelude_token_len and len(chunk_paras) > 0:
        if prev_len <= len(doc_dct['paragraphs']): # skip previously processed chunks
            chunk = '.'.join(chunk_paras)
            para_dct = {}
            para_dct['paragraph_text'] = chunk
            try:
                embedding = vectorizer([f"{prelude}{chunk}"])
                eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                para_dct['embedding'] = eem
            except Exception as ex:
                print(f"Exception {ex} while creating residual para embedding")
            doc_dct['paragraphs'].append(para_dct)
    return doc_dct

def read_html(email, filename:str, fileid:str, file_io:io.BytesIO, mtime:datetime.datetime, prev_paras) -> Dict[str, Union[str,Dict[str, str]]]:
    vectorizer = get_vectorizer(email, [])
    doc_dct={"filename": filename, "fileid": fileid, "mtime": mtime, "paragraphs": prev_paras}
    prev_len = len(prev_paras)
    html_content = file_io.getvalue().decode('utf-8')
    soup = BeautifulSoup(html_content, 'html.parser')
    pgs = soup.find_all('p')
    text_paragraphs = []
    for para in pgs:
        text = para.get_text()
        nonl_text = text.replace('\n', ' ').replace('\r', '')
        nonl_ss = re.sub(' +', ' ', nonl_text)
        text_paragraphs.append(nonl_ss)
    prelude = f"The filename is {filename} and the paragraphs are:"
    prelude_token_len = vectorizer.get_token_count(prelude)
    chunk_len = prelude_token_len
    chunk_paras = []
    for para in text_paragraphs:
        if para:
            para_len = vectorizer.get_token_count(para)
            if para_len + chunk_len >= 512:
                if prev_len > len(doc_dct['paragraphs']): # skip previously processed chunks
                    chunk_len = prelude_token_len
                    chunk_paras = []
                    continue
                chunk = '.'.join(chunk_paras)
                para_dct = {'paragraph_text': chunk}
                try:
                    embedding = vectorizer([f"{prelude}{chunk}"])
                    eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                    para_dct['embedding'] = eem
                except Exception as ex:
                    print(f"Exception {ex} while creating para embedding")
                doc_dct['paragraphs'].append(para_dct)
                chunk_len = prelude_token_len
                chunk_paras = []
                if lambda_timelimit_exceeded():
                    doc_dct['partial'] = "true"
                    print(f"read_html: Lambda timelimit exceeded when reading docx file. Breaking..")
                    break
            else:
                chunk_paras.append(para)
                chunk_len += para_len
    if chunk_len > prelude_token_len and len(chunk_paras) > 0:
        if prev_len <= len(doc_dct['paragraphs']): # skip previously processed chunks
            chunk = '.'.join(chunk_paras)
            para_dct = {}
            para_dct['paragraph_text'] = chunk
            try:
                embedding = vectorizer([f"{prelude}{chunk}"])
                eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                para_dct['embedding'] = eem
            except Exception as ex:
                print(f"Exception {ex} while creating residual para embedding")
            doc_dct['paragraphs'].append(para_dct)
    return doc_dct

def process_docx(email, file_item, filename, fileid, bio):
    if 'partial' in file_item and 'paragraphs' in file_item:
        del file_item['partial']
        prev_paras = file_item['paragraphs']
        print(f"process_docx: fn={filename}. found partial. len(prev_paras)={len(prev_paras)}")
    else:
        prev_paras = []
        print(f"process_docx: fn={filename}. did not find partial")
    doc_dict = read_docx(email, filename, fileid, bio, file_item['mtime'], prev_paras)
    file_item['paragraphs'] = doc_dict['paragraphs']
    file_item['filetype'] = 'docx'
    if 'partial' in doc_dict:
        file_item['partial'] = 'true'

def process_txt(email, file_item, filename, fileid, bio):
    if 'partial' in file_item and 'paragraphs' in file_item:
        del file_item['partial']
        prev_paras = file_item['paragraphs']
        print(f"process_txt: fn={filename}. found partial. len(prev_paras)={len(prev_paras)}")
    else:
        prev_paras = []
        print(f"process_txt: fn={filename}. did not find partial")
    doc_dict = read_txt(email, filename, fileid, bio, file_item['mtime'], prev_paras)
    file_item['paragraphs'] = doc_dict['paragraphs']
    file_item['filetype'] = 'txt'
    if 'partial' in doc_dict:
        file_item['partial'] = 'true'

def process_html(email, file_item, filename, fileid, bio):
    if 'partial' in file_item and 'paragraphs' in file_item:
        del file_item['partial']
        prev_paras = file_item['paragraphs']
        print(f"process_html: fn={filename}. found partial. len(prev_paras)={len(prev_paras)}")
    else:
        prev_paras = []
        print(f"process_html: fn={filename}. did not find partial")
    doc_dict = read_html(email, filename, fileid, bio, file_item['mtime'], prev_paras)
    file_item['paragraphs'] = doc_dict['paragraphs']
    file_item['filetype'] = 'html'
    if 'partial' in doc_dict:
        file_item['partial'] = 'true'

def process_gh_issues_zip(email, file_item, filename, fileid, mimetype, zip_path):
    vectorizer = get_vectorizer(email, [])
    # Each issue is one paragraph
    if 'partial' in file_item and 'paragraphs' in file_item:
        del file_item['partial']
        prev_paras = file_item['paragraphs']
        print(f"process_gh_issues_zip: fn={filename}. found partial. len(prev_paras)={len(prev_paras)}")
    else:
        prev_paras = []
        print(f"process_gh_issues_zip: fn={filename}. did not find partial")
    timelimit_exceeded = False
    # Create a temporary directory
    with tempfile.TemporaryDirectory() as tmpdirname:
        print('Created temporary directory:', tmpdirname)
        # Unzip the downloaded file
        with zipfile.ZipFile(zip_path, 'r') as zip_ref:
            zip_ref.extractall(tmpdirname)
            print(f"Files extracted to {tmpdirname}")

        extracted_files = os.listdir(tmpdirname)
        issues_processed = 0
        for jsonfile in extracted_files:
            try:
                if not jsonfile.endswith(".json"):
                    continue
                issues_processed += 1
                htmlfile = f"{jsonfile[:-5]}.html"
                print(f"process_gh_issues_zip: Processed {issues_processed} issues. Now processing {jsonfile}")
                with open(os.path.join(tmpdirname, jsonfile), 'r') as jfp:
                    js = json.load(jfp)
                with open(os.path.join(tmpdirname, htmlfile), 'r') as hfp:
                    html_content = hfp.read()
                    soup = BeautifulSoup(html_content, 'html.parser')
                    pgs = soup.find_all('p')
                    text_paragraphs = []
                    for para in pgs:
                        text = para.get_text()
                        nonl_text = text.replace('\n', ' ').replace('\r', '')
                        nonl_ss = re.sub(' +', ' ', nonl_text)
                        text_paragraphs.append(nonl_ss)
                txt = f"bug id {js['number']} is {js['title']} and it is in state {js['state']}. Details of the issue are as follows: {'.'.join(text_paragraphs)}"
                print(f"process_gh_issues_zip: txt={txt}")
                para_dct = {}
                para_dct['paragraph_text'] = txt
                para_dct['html_url'] = js['html_url']
                embedding = vectorizer([txt])
                eem = base64.b64encode(pickle.dumps(embedding)).decode('ascii')
                para_dct['embedding'] = eem
                para_num = js['number'] -1 # github issues are 1 based while paragraphs are 0 based
                print(f"process_gh_issues_zip: issue#{js['number']}, para_num={para_num}, len(prev_paras)={len(prev_paras)}")
                if para_num < len(prev_paras):
                    # updating existing issue
                    print(f"process_gh_issues_zip: Updating existing")
                    prev_paras[para_num] = para_dct
                elif len(prev_paras) == para_num:
                    # next issue
                    print(f"process_gh_issues_zip: appending next")
                    prev_paras.append(para_dct)
                else:
                    # sparse. Fill in with Nones
                    print(f"process_gh_issues_zip: sparse filling. extending {para_num - len(prev_paras)} with None")
                    for i in range(para_num - len(prev_paras)):
                        prev_paras.append(None)
                    prev_paras.append(para_dct)
                if lambda_timelimit_exceeded():
                    print(f"process_gh_issues_zip: Processed {issues_processed} issues. Exceeded lambda timelimit. Trying to extend..")
                    extend_lock_time(email, index_dir, lambda_time_left_seconds())
                    if lambda_timelimit_exceeded():
                        print(f"process_gh_issues_zip: Processed {issues_processed} issues. Extending timelimit appears to not have worked. Breaking..")
                        timelimit_exceeded = True
                        break
            except Exception as ex:
                print(f"process_gh_issues_zip: Caught {ex} while processing {jsonfile}. Ignoring this issue and carrying on..")
    file_item['paragraphs'] = prev_paras
    file_item['filetype'] = 'gh-issues.zip'
    if timelimit_exceeded:
        file_item['partial'] = "true"

def process_pptx(email, file_item, filename, fileid, bio):
    if 'partial' in file_item and 'slides' in file_item:
        del file_item['partial']
        prev_slides = file_item['slides']
        print(f"process_pptx: fn={filename}. found partial. len(prev_slides)={len(prev_slides)}")
    else:
        prev_slides = []
        print(f"process_pptx: fn={filename}. did not find partial")
    ppt = read_pptx(email, filename, fileid, bio, file_item['mtime'], prev_slides)
    file_item['slides'] = ppt['slides']
    file_item['filetype'] = 'pptx'
    if 'partial' in ppt:
        file_item['partial'] = 'true'

def process_xlsx(email, file_item, filename, fileid, bio):
    if 'partial' in file_item and 'rows' in file_item:
        del file_item['partial']
        prev_rows = file_item['rows']
        print(f"process_xlsx: fn={filename}. found partial. len(prev_workbook)={len(prev_workbook)}")
    else:
        prev_rows = []
        print(f"process_xlsx: fn={filename}. did not find partial")
    xlsx = read_xlsx(email, filename, fileid, bio, file_item['mtime'], prev_rows)
    file_item['rows'] = xlsx['rows']
    file_item['filetype'] = 'xlsx'
    if 'partial' in xlsx:
        file_item['partial'] = 'true'

class StorageReader:
    def read(self, fileid, filename, mimetype):
        pass
    def descriptor(self):
        pass
    def download_to_local(self, fileid, filename, mimetype, localfn):
        pass

class LocalFileReader(StorageReader):
    def __init__(self, docs_dir, file_details, dir_details):
        self._docs_dir = docs_dir
        self._file_details = file_details
        self._dir_details = dir_details

    def read(self, fileid, filename, mimetype):
        print(f"LocalFileReader.read: Entered. fileid={fileid}, filename={filename}, mimetype={mimetype}")
        finfo = self._file_details[fileid]
        path = ""
        while True:
            if not 'parents' in finfo:
                break
            fileid = finfo['parents'][0]
            dirinfo = self._dir_details[fileid]
            path = os.path.join(dirinfo['filename'], path)
            finfo = dirinfo
        full_fn = os.path.join(self._docs_dir, path, filename)
        print(f"LocalFileReader._read: full_fn={full_fn}")
        with open(full_fn, 'rb') as rfp:
            return io.BytesIO(rfp.read())

    def download_to_local(self, fileid, filename, mimetype, localfn):
        bio = self.read(fileid, filename, mimetype)
        with open(localfn, 'wb') as wfp:
            wfp.write(bio.getvalue())

    def descriptor(self):
        return 'local'

class GoogleDriveReader(StorageReader):
    def __init__(self, service, access_token):
        self._service = service
        self._access_token = access_token
    def read(self, fileid, filename, mimetype):
        print(f"GoogleDriveReader.read: Entered. fileid={fileid}, filename={filename}, mimetype={mimetype}")
        if mimetype == 'application/vnd.google-apps.presentation':
            print(f"GoogleDriveReader.read: fileid={fileid} calling export1")
            return export_gdrive_file(self._access_token, fileid, 'pptx')
        elif mimetype == 'application/vnd.google-apps.document':
            print(f"GoogleDriveReader.read: fileid={fileid} calling export2")
            return export_gdrive_file(self._access_token, fileid, 'docx')
        elif mimetype == 'application/vnd.google-apps.spreadsheet':
            print(f"GoogleDriveReader.read: fileid={fileid} calling export3")
            return export_gdrive_file(self._access_token, fileid, 'xlsx')
        else:
            print(f"GoogleDriveReader.read: fileid={fileid} calling download")
            return download_gdrive_file(self._service, fileid, filename)
    def download_to_local(self, fileid, filename, mimetype, localfn):
        print(f"GogleDriveReader.download_to_local: Entered. fileid={fileid}, filename={filename}, localfn={localfn}")
        request = self._service.files().get_media(fileId=fileid)
        fh = io.FileIO(localfn, 'w+b')
        # Download the file
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        while not done:
            status, done = downloader.next_chunk()
            print(f"Download {int(status.progress() * 100)}%.")
        # Close the file handler after the download
        fh.close()
        return

    def descriptor(self):
        return 'gdrive'

class DropboxReader(StorageReader):
    def __init__(self, token):
        self._token = token
    def read(self, fileid, filename, mimetype):
        print(f"DropboxReader.read: Entered. fileid={fileid}, filename={filename}")
        url = 'https://content.dropboxapi.com/2/files/download'
        headers = {'Authorization': f"Bearer {self._token}", 'Dropbox-API-Arg': json.dumps({'path': fileid})}
        r = requests.get(url, stream=True, headers=headers)
        file = io.BytesIO()
        for chunk in r.iter_content(chunk_size=1024): 
            if chunk:
                file.write(chunk)
        file.seek(0, os.SEEK_SET)
        return file
    def download_to_local(self, fileid, filename, mimetype, localfn):
        print(f"DropboxReader.download_to_local: Entered. fileid={fileid}, filename={filename}, localfn={localfn}")
        with open(localfn, 'w+b') as wfp:
            url = 'https://content.dropboxapi.com/2/files/download'
            headers = {'Authorization': f"Bearer {self._token}", 'Dropbox-API-Arg': json.dumps({'path': fileid})}
            r = requests.get(url, stream=True, headers=headers)
            file = io.BytesIO()
            for chunk in r.iter_content(chunk_size=1024): 
                if chunk:
                    wfp.write(chunk)
            fh.close()
            return
    def descriptor(self):
        return 'dropbox'

def _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, force):
    now = int(time.time())
    if not force:
        if prev_update and (now - prev_update) < 60:
            return False, 0
    len_unmodified = 0
    unmodified_size = 0
    if unmodified:
        len_unmodified = len(unmodified)
        for ue_key, ue_value in unmodified.items():
            if 'size' in ue_value:
                unmodified_size = unmodified_size + int(ue_value['size'])
    len_needs_embedding = 0
    needs_embedding_size = 0
    if needs_embedding:
        len_needs_embedding = len(needs_embedding)
        for ne_key, ne_value in needs_embedding.items():
            if 'size' in ne_value:
                needs_embedding_size = needs_embedding_size + int(ne_value['size'])
    len_done_embedding = 0
    done_embedding_size = 0
    if done_embedding:
        len_done_embedding = len(done_embedding)
        for de_key, de_value in done_embedding.items():
            if 'size' in de_value:
                done_embedding_size = done_embedding_size + int(de_value['size'])
    print(f"_update_progress_file: {storage_reader.descriptor()} unmodified_size={unmodified_size}, needs_embedding_size={needs_embedding_size}, done_embedding_size={done_embedding_size}")
    _write_progress_file(index_dir, s3client, bucket, prefix, storage_reader.descriptor(),
                len_unmodified + len_done_embedding, unmodified_size + done_embedding_size,
                len_needs_embedding - len_done_embedding, needs_embedding_size - done_embedding_size)
    return True, now

def process_files(email, storage_reader:StorageReader, unmodified, needs_embedding, index_dir, s3client, bucket, prefix) -> Tuple[Dict[str, Dict[str, Any]], bool] : 
    """ processs the files in google drive. returns a dict:  { fileid: {filename, fileid, mtime} }
    'service': google drive service ; 'needs_embedding':  the list of files as a dict that needs to be embedded; s3client, bucket, prefix: location of the vector db/index file 
    Note: will stop processing files after PERIOIDIC_PROCESS_FILES_TIME_LIMIT minutes (env variable)
    """
    print(f"process_files: Entered")
    time_limit_exceeded = False
    prev_update = 0
    done_embedding = {}
    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
    if status:
        prev_update = updtime
    to_del_from_needs_embedding = []
    for fileid, file_item in needs_embedding.items():
        filename = file_item['filename']
        path = file_item['path'] if 'path' in file_item else ''
        mimetype = file_item['mimetype'] if 'mimetype' in file_item else ''
        # handle errors like these: a docx file that was deleted (in Trash) but is still visible in google drive api: raise BadZipFile("File is not a zip file")
        try:
            if 'size' in file_item and int(file_item['size']) == 0: 
                if not (mimetype == 'application/vnd.google-apps.presentation' \
                        or mimetype == 'application/vnd.google-apps.document' \
                        or mimetype == 'application/vnd.google-apps.spreadsheet'):
                    print(f"skipping google drive with file size == 0, and not of type google sheet, slides or docs: {file_item}")
                    to_del_from_needs_embedding.append(fileid)
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
                    continue
            if filename.lower().endswith('.pptx'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                process_pptx(email, file_item, filename, fileid, bio)
                done_embedding[fileid] = file_item
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            elif filename.lower().endswith('.docx'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                process_docx(email, file_item, filename, fileid, bio)
                done_embedding[fileid] = file_item
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            elif filename.lower().endswith('.xlsx'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                process_xlsx(email, file_item, filename, fileid, bio)
                done_embedding[fileid] = file_item
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            elif filename.lower().endswith('.doc'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                tfd, tfn = tempfile.mkstemp(suffix=".doc", dir="/tmp")
                with os.fdopen(tfd, "wb") as wfp:
                    wfp.write(bio.getvalue())
                bio.close()
                rv = subprocess.run(['/opt/libreoffice7.6/program/soffice', '--headless', '--convert-to', 'docx:MS Word 2007 XML', tfn, '--outdir', '/tmp'],
                                        stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
                for line in rv.stdout.splitlines():
                    print(line.decode('utf-8'))
                os.remove(tfn)
                if rv.returncode == 0:
                    print(f"Successfully converted {filename} to docx. temp file is {tfn}, Proceeding with embedding generation")
                    with open(f"{tfn}x", 'rb') as fp:
                        process_docx(email, file_item, filename, fileid, fp)
                        done_embedding[fileid] = file_item
                        status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                        if status:
                            prev_update = updtime
                    os.remove(f'{tfn}x')
                else:
                    print(f"Failed to convert {filename} to docx. Return value {rv.returncode}. Not generating embeddings")
            elif filename.lower().endswith('.ppt'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                tfd, tfn = tempfile.mkstemp(suffix=".ppt", dir="/tmp")
                with os.fdopen(tfd, "wb") as wfp:
                    wfp.write(bio.getvalue())
                bio.close()
                rv = subprocess.run(['/opt/libreoffice7.6/program/soffice', '--headless', '--convert-to', 'pptx', tfn, '--outdir', '/tmp'],
                                    stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
                for line in rv.stdout.splitlines():
                    print(line.decode('utf-8'))
                os.remove(tfn)
                if rv.returncode == 0:
                    print(f"Successfully converted {filename} to pptx. temp file is {tfn}, Proceeding with embedding generation")
                    with open(f"{tfn}x", 'rb') as fp:
                        process_pptx(email, file_item, filename, fileid, fp)
                        done_embedding[fileid] = file_item
                        status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                        if status:
                            prev_update = updtime
                    os.remove(f'{tfn}x')
                else:
                    print(f"Failed to convert {filename} to pptx. Return value {rv.returncode}. Not generating embeddings")
            elif (filename.lower().endswith('.pdf') or mimetype == 'application/pdf'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                if 'partial' in file_item and 'paragraphs' in file_item:
                    del file_item['partial']
                    prev_paras = file_item['paragraphs']
                    print(f"process_files: fn={filename}. found partial. len(prev_paras)={len(prev_paras)}")
                else:
                    prev_paras = []
                    print(f"process_files: fn={filename}. did not find partial")
                vectorizer = get_vectorizer(email, [])
                pdf_dict:Dict[str,Any] = read_pdf(email, index_dir, filename, fileid, bio, file_item['mtime'], vectorizer, prev_paras)
                file_item['paragraphs'] = pdf_dict['paragraphs']
                file_item['filetype'] = 'pdf'
                if 'partial' in pdf_dict:
                    file_item['partial'] = 'true'
                done_embedding[fileid] = file_item                    
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            elif mimetype == 'application/vnd.google-apps.presentation':
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                if not bio:
                    print(f"process_files: Unable to export {filename} to pptx format")
                    to_del_from_needs_embedding.append(fileid)
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
                else:
                    process_pptx(email, file_item, filename, fileid, bio)
                    done_embedding[fileid] = file_item
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
            elif mimetype == 'application/vnd.google-apps.document':
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                if not bio:
                    print(f"process_files: Unable to export {filename} to docx format")
                    to_del_from_needs_embedding.append(fileid)
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
                else:
                    process_docx(email, file_item, filename, fileid, bio)
                    done_embedding[fileid] = file_item
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
            elif mimetype == 'application/vnd.google-apps.spreadsheet':
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                if not bio:
                    print(f"process_files: Unable to export {filename} to xlsx format")
                    to_del_from_needs_embedding.append(fileid)
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
                else:
                    process_xlsx(email, file_item, filename, fileid, bio)
                    done_embedding[fileid] = file_item
                    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                    if status:
                        prev_update = updtime
            elif (filename.lower().endswith('.txt') or mimetype == 'text/plain'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                process_txt(email, file_item, filename, fileid, bio)
                done_embedding[fileid] = file_item
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            elif (filename.lower().endswith('.html') or mimetype == 'text/html'):
                bio:io.BytesIO = storage_reader.read(fileid, filename, mimetype)
                process_html(email, file_item, filename, fileid, bio)
                done_embedding[fileid] = file_item
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            elif filename.lower().endswith('.gh-issues.zip'):
                storage_reader.download_to_local(fileid, filename, mimetype, '/tmp/downloaded_file.zip')
                process_gh_issues_zip(email, file_item, filename, fileid, mimetype, '/tmp/downloaded_file.zip')
                done_embedding[fileid] = file_item
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
            else:
                print(f"process_files: skipping unknown file type {filename} mimetype={mimetype}")
                to_del_from_needs_embedding.append(fileid)
                status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
                if status:
                    prev_update = updtime
        except Exception as e:
            print(f"process_files(): skipping filename={filename} with fileid={fileid} due to exception={e}")
            traceback.print_exc()
            to_del_from_needs_embedding.append(fileid)
            status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, prev_update, index_dir, s3client, bucket, prefix, False)
            if status:
                prev_update = updtime
        if lambda_timelimit_exceeded():
            print(f"process_files: lambda_timelimit_exceeded. Breaking..")
            time_limit_exceeded = True
            break
    # final update
    for itm in to_del_from_needs_embedding:
        del needs_embedding[itm]
    status, updtime = _update_progress_file(storage_reader, unmodified, needs_embedding, done_embedding, 0, index_dir, s3client, bucket, prefix, True)
    if status:
        prev_update = updtime
    return done_embedding, time_limit_exceeded

def _create_index_metadata_json_local(in_index_metadata:IndexMetadata):
    with open(INDEX_METADATA_JSON, "w") as f:
         json.dump(jsons.dump(in_index_metadata), f)
         print(f"{INDEX_METADATA_JSON} after create = {in_index_metadata}")

def _update_index_metadata_json_local(in_index_metadata:IndexMetadata):
    # update the metadata file
    with open(INDEX_METADATA_JSON, "r") as f:
        index_metadata = jsons.load(json.load(f),IndexMetadata)
        print(f"{INDEX_METADATA_JSON} before update = {index_metadata}")

    if in_index_metadata.jsonl_last_modified: index_metadata.jsonl_last_modified = in_index_metadata.jsonl_last_modified
    if in_index_metadata.index_flat_last_modified: index_metadata.index_flat_last_modified = in_index_metadata.index_flat_last_modified
    if in_index_metadata.index_ivfadc_last_modified: index_metadata.index_ivfadc_last_modified = in_index_metadata.index_ivfadc_last_modified
    if in_index_metadata.bm25s_index_last_modified: index_metadata.bm25s_index_last_modified = in_index_metadata.bm25s_index_last_modified
    
    with open(INDEX_METADATA_JSON, "w") as f:
         json.dump(jsons.dump(index_metadata), f)
         print(f"{INDEX_METADATA_JSON} after  update = {index_metadata}")

def _read_index_metadata_json_local() -> IndexMetadata:
    index_metadata:IndexMetadata = None
    with open(INDEX_METADATA_JSON, "r") as f:
        index_metadata = jsons.load(json.load(f),IndexMetadata)
        print(f"read {INDEX_METADATA_JSON} = {index_metadata}")
    return index_metadata
    
def update_files_index_jsonl(done_embedding, unmodified, index_dir, bucket, user_prefix, s3client):
    print(f"update_files_index_jsonl: Entered index_dir={index_dir}, bucket={bucket}, use_prefilx={user_prefix} . num new embeddings={len(done_embedding.items())}, num unmodified={len(unmodified.items())}")
    # consolidate unmodified and done_embedding
    for fileid, file_item in done_embedding.items():
        unmodified[fileid] = file_item
    unsorted_all_files = []
    for fileid, file_item in unmodified.items():
        unsorted_all_files.append(file_item)
    sorted_all_files = sorted(unsorted_all_files, key=lambda x: x['mtime'], reverse=True)
    files_index_jsonl_gz = FILES_INDEX_JSONL_GZ
    with gzip.open(files_index_jsonl_gz, "wb") as wfp:
        for file_item in sorted_all_files:
            file_item['mtime'] = to_rfc3339(file_item['mtime'])
            wfp.write((json.dumps(file_item) + "\n").encode())
    
    if index_dir:
        print(f"Copying  {files_index_jsonl_gz}  Size of file={os.path.getsize(files_index_jsonl_gz)} to index_dir {index_dir}")
        shutil.copy2(files_index_jsonl_gz, os.path.join(index_dir, os.path.basename(files_index_jsonl_gz)))
    else:
        print(f"Uploading  {files_index_jsonl_gz}  Size of file={os.path.getsize(files_index_jsonl_gz)}")
        s3client.upload_file(files_index_jsonl_gz, bucket, f"{user_prefix}/{os.path.basename(files_index_jsonl_gz)}")
    os.remove(FILES_INDEX_JSONL_GZ)
    # remove stale files_index.jsonl
    if os.path.exists(FILES_INDEX_JSONL): os.remove(FILES_INDEX_JSONL)
    
    # create an empty index_metadata_json if it doesn't already exist: this is needed when the index is first being created and none of the files files_index.jsonl.gz, index_metadata.json, faiss_index_flat and faiss_index_ivfadc exist.
    if not os.path.exists(INDEX_METADATA_JSON):
        _create_index_metadata_json_local(IndexMetadata(jsonl_last_modified=time.time()))
    else:
        # update index_metadata_json
        _update_index_metadata_json_local(IndexMetadata(jsonl_last_modified=time.time()))
    if index_dir:
        print(f"Copying  {INDEX_METADATA_JSON}  Size of file={os.path.getsize(INDEX_METADATA_JSON)} to index_dir {index_dir}")
        shutil.copy2(INDEX_METADATA_JSON, os.path.join(index_dir, os.path.basename(INDEX_METADATA_JSON)))
    else:
        print(f"Uploading  {INDEX_METADATA_JSON}  Size of file={os.path.getsize(INDEX_METADATA_JSON)}")
        s3client.upload_file(INDEX_METADATA_JSON, bucket, f"{user_prefix}/{os.path.basename(INDEX_METADATA_JSON)}")

def _delete_faiss_index(email:str, index_dir, s3client:S3Client, bucket:str, prefix:str, sub_prefix:str, user_prefix:str ):
    """  user_prefix == prefix + email + sub_prefix """
    for faiss_index_fname in (FAISS_INDEX_FLAT, FAISS_INDEX_IVFADC, BM25S_INDEX_TAR):
        faiss_s3_key = f"{user_prefix}/{os.path.basename(faiss_index_fname)}"
        if os.path.exists(faiss_index_fname): 
            print(f"Deleting faiss index {faiss_index_fname} in local filesystem")
            os.remove(faiss_index_fname)
        try:
            if index_dir:
                faiss_fbase = os.path.join(index_dir, os.path.basename(faiss_index_fname))
                print(f"Deleting faiss index {faiss_fbase}")
                if os.path.exists(faiss_fbase):
                    os.remove(faiss_fbase)
            else:
                s3client.head_object(Bucket=bucket, Key=faiss_s3_key)
                print(f"Deleting faiss index {faiss_index_fname} in S3 key={faiss_s3_key}")
                s3client.delete_object(Bucket=bucket, Key=faiss_s3_key)
        except Exception as e:
            pass

def create_tar_file(directory_path, tar_file_path):
    with tarfile.open(tar_file_path, "w") as tar:
        tar.add(directory_path, arcname=os.path.basename(directory_path))

def build_and_save_faiss(email, index_dir, s3client, bucket, prefix, sub_prefix, user_prefix, doc_storage_type:DocStorageType):
    """ Note that user_prefix == prefix + email + sub_prefix """
    # now build the index.
    try:
        faiss_rm:FaissRM = init_vdb(email, index_dir, s3client, bucket, prefix,
                                doc_storage_type=doc_storage_type,
                                sub_prefix=sub_prefix, tracebuf=[])
    except Exception as ex:
        print(f"build_and_save_faiss: Caught {ex}")
        traceback.print_exc()

    # save the created index
    faiss_flat_fname = FAISS_INDEX_FLAT
    faiss.write_index(faiss_rm.get_index_flat(), faiss_flat_fname)
    if index_dir:
        print(f"Copying faiss flat index {faiss_flat_fname}  Size of index={os.path.getsize(faiss_flat_fname)} to index_dir={index_dir}")
        shutil.copy2(faiss_flat_fname, os.path.join(index_dir, os.path.basename(faiss_flat_fname)))
    else:
        print(f"Uploading faiss flat index {faiss_flat_fname}  Size of index={os.path.getsize(faiss_flat_fname)}")
        s3client.upload_file(faiss_flat_fname, bucket, f"{user_prefix}/{os.path.basename(faiss_flat_fname)}")
    os.remove(faiss_flat_fname)
    
    # update index_metadata_json    
    _update_index_metadata_json_local(IndexMetadata(index_flat_last_modified=time.time()))
    if index_dir:
        print(f"Copying faiss flat index {INDEX_METADATA_JSON}  Size of index={os.path.getsize(INDEX_METADATA_JSON)} to index_dir={index_dir}")
        shutil.copy2(INDEX_METADATA_JSON, os.path.join(index_dir, os.path.basename(INDEX_METADATA_JSON)))
    else:
        print(f"Uploading  {INDEX_METADATA_JSON}  Size of file={os.path.getsize(INDEX_METADATA_JSON)}")
        s3client.upload_file(INDEX_METADATA_JSON, bucket, f"{user_prefix}/{os.path.basename(INDEX_METADATA_JSON)}")
    
    # save the created index
    faiss_ivfadc_fname = FAISS_INDEX_IVFADC
    faiss.write_index(faiss_rm.get_index_ivfadc(), faiss_ivfadc_fname)
    if index_dir:
        print(f"Copying faiss flat index {faiss_ivfadc_fname}  Size of index={os.path.getsize(faiss_ivfadc_fname)} to index_dir={index_dir}")
        shutil.copy2(faiss_ivfadc_fname, os.path.join(index_dir, os.path.basename(faiss_ivfadc_fname)))
    else:
        print(f"Uploading faiss ivfadc index {faiss_ivfadc_fname}.  Size of index={os.path.getsize(faiss_ivfadc_fname)}")
        s3client.upload_file(faiss_ivfadc_fname, bucket, f"{user_prefix}/{os.path.basename(faiss_ivfadc_fname)}")
    os.remove(faiss_ivfadc_fname)

    # save the bm25s retriever
    bm25s_retriever = faiss_rm.get_bm25s_retriever()
    bm25s_retriever.save(BM25S_INDEX)
    create_tar_file(BM25S_INDEX, BM25S_INDEX_TAR)
    _update_index_metadata_json_local(IndexMetadata(bm25s_index_last_modified=time.time()))
    if index_dir:
        print(f"Copying faiss flat index {BM25S_INDEX_TAR}  Size of index={os.path.getsize(BM25S_INDEX_TAR)} to index_dir={index_dir}")
        shutil.copy2(BM25S_INDEX_TAR, os.path.join(index_dir, os.path.basename(BM25S_INDEX_TAR)))
    else:
        print(f"Uploading bm25s index {BM25S_INDEX_TAR}.  Size of index={os.path.getsize(BM25S_INDEX_TAR)}")
        s3client.upload_file(BM25S_INDEX_TAR, bucket, f"{user_prefix}/{os.path.basename(BM25S_INDEX_TAR)}")
    os.remove(BM25S_INDEX_TAR)

    # update index_metadata_json
    _update_index_metadata_json_local(IndexMetadata(index_ivfadc_last_modified=time.time()))
    if index_dir:
        print(f"Copying faiss flat index {INDEX_METADATA_JSON}  Size of index={os.path.getsize(INDEX_METADATA_JSON)} to index_dir={index_dir}")
        shutil.copy2(INDEX_METADATA_JSON, os.path.join(index_dir, os.path.basename(INDEX_METADATA_JSON)))
    else:
        print(f"Uploading  {INDEX_METADATA_JSON}  Size of file={os.path.getsize(INDEX_METADATA_JSON)}")
        s3client.upload_file(INDEX_METADATA_JSON, bucket, f"{user_prefix}/{os.path.basename(INDEX_METADATA_JSON)}")

def _build_and_save_faiss_if_needed(email:str, index_dir, s3client, bucket:str, prefix:str, sub_prefix:str, user_prefix:str, doc_storage_type:DocStorageType ):
    """  user_prefix == prefix + email + sub_prefix 
    """
    
    # at this point, we must have the index_metadata.json file: it must have been downloaded above (and is unmodified as no files need to be embedded) or it was modified above (as files needed to be embedded)
    index_metadata:IndexMetadata = _read_index_metadata_json_local()
    # if faiss index needs to be updated and we have at least 5 minutes
    if index_metadata.is_vdb_index_stale():
        time_left:int = lambda_time_left_seconds()
        extend_lock_time(email, index_dir, time_left)
        if time_left > 300:
            print(f"updating faiss index for s3://{bucket}/{user_prefix} as time left={time_left} > 300 seconds")
            build_and_save_faiss(email, index_dir, s3client, bucket, prefix, sub_prefix, user_prefix, doc_storage_type)
        else:
            print(f"Not updating faiss index for s3://{bucket}/{user_prefix} as time left={time_left} < 300 seconds.  Deleting stale faiss indexes if needed..")
            _delete_faiss_index(email, index_dir, s3client, bucket, prefix, sub_prefix, user_prefix)
            
    else:
        print(f"Not updating faiss index for s3://{bucket}/{user_prefix} as it isn't stale: {index_metadata}")

# Returns entries, dropbox_next_page_token
def _get_dropbox_listing(dbx, dropbox_next_page_token):
    print(f"_get_dropbox_listing: Entered dropbox_next_page_token={dropbox_next_page_token}")
    try: 
        all_entries = []
        cursor = dropbox_next_page_token
        has_more = True
        while has_more:
            if not cursor:
                res = dbx.files_list_folder("", recursive=True)
            else:
                res = dbx.files_list_folder_continue(cursor=cursor)
            all_entries.extend(res.entries)
            has_more = res.has_more
            cursor = res.cursor
        return all_entries, cursor
    except Exception as e: 
        print(str(e)) 
        return None

# returns unmodified, needs_embedding
def _calc_filelists_dropbox(entries, s3_index):
    needs_embedding = {}
    num_deleted_files = 0
    for entry in entries:
        if isinstance(entry, dropbox.files.FolderMetadata):
            print(f"_calc_filelists_dropbox: type=DIR id={entry.id} path_display={entry.path_display}, name={entry.name}")
        elif isinstance(entry, dropbox.files.FileMetadata):
            print(f"_calc_filelists_dropbox: type=FILE id={entry.id} path_display={entry.path_display}, name={entry.name}, mtime={entry.server_modified}")
            mtime = entry.server_modified.replace(tzinfo=timezone.utc)
            needs_embedding[entry.id] = {'filename': entry.name,
                                        'fileid': entry.id,
                                        'path': entry.path_display,
                                        'mtime': mtime}
        elif isinstance(entry, dropbox.files.DeletedMetadata):
            print(f"_calc_filelists_dropbox: type=DELETED path_display={entry.path_display}, name={entry.name}")
            to_delete = []
            for fileid, s3entry in s3_index.items():
                if s3entry['path'].startswith(entry.path_display):
                    print(f"_calc_filelists_dropbox: DELETED Yes! deleted={entry.path_display}. s3entry={s3entry['path']}")
                    to_delete.append(fileid)
                else:
                    print(f"_calc_filelists_dropbox: DELETED No! deleted={entry.path_display}. s3entry={s3entry['path']}")
            for td in to_delete:
                del s3_index[td]
                if td in needs_embedding:
                    del needs_embedding[td]
            num_deleted_files += len(to_delete)
        else:
            print(f"_calc_filelists_dropbox: Unknown entry type {entry}. Ignoring and continuing..")
    # if there is any partial entry in s3_index, move it to needs_embedding
    to_move = []
    for fid, entry in s3_index.items():
        if 'partial' in entry:
            to_move.append(fid)
    for tm in to_move:
        needs_embedding[tm] = s3_index[tm]
        del s3_index[tm]
    print(f"_calc_filelists_dropbox: return. len s3_index={len(s3_index.items())}, len needs_embedding={len(needs_embedding.items())}, num_deleted_files={num_deleted_files}")
    return s3_index, needs_embedding, num_deleted_files

def partial_present(done_embedding):
    for fileid, file_item in done_embedding.items():
        if 'partial' in file_item:
            return True
    return False

def update_index_for_user_dropbox(email, s3client, bucket:str, prefix:str, dropbox_next_page_token=None):
    print(f'update_index_for_user_dropbox: Entered. {email}')
    item = get_user_table_entry(email)
    sub_prefix = "dropbox"
    user_prefix = f"{prefix}/{email}/{sub_prefix}"
    try:
        s3_index = get_s3_index(None, s3client, bucket, user_prefix)
        access_token = refresh_user_dropbox(item)
    except HttpError as error:
        print(f"HttpError occurred: {error}")
        return
    except Exception as ex:
        print(f"Exception occurred: {ex}")
        return
    if not access_token:
        print(f"Error. Failed to create access token. Not updating dropbox index")
        return

    try: 
        dbx = dropbox.Dropbox(access_token) 
        entries, dropbox_next_page_token = _get_dropbox_listing(dbx, dropbox_next_page_token)
        unmodified, needs_embedding, num_deleted_files = _calc_filelists_dropbox(entries, s3_index)
    except Exception as ex:
        print(f"Exception occurred: {ex}")

    # Move items in 'partial' state from umodified to needs_embedding
    to_del=[]
    for fid, entry in unmodified.items():
        if 'partial' in entry:
            pth = entry['path'] if 'path' in entry else ''
            print(f"update.dropbox: fileid={fid}, path={pth}, filename={entry['filename']} is in partial state. Moving from unmodified to needs_embedding")
            needs_embedding[fid] = entry
            to_del.append(fid)
    for td in to_del:
        del unmodified[td]

    done_embedding, time_limit_exceeded = process_files(email, DropboxReader(access_token), unmodified, needs_embedding, None, s3client, bucket, user_prefix)
    if len(done_embedding.items()) > 0 or num_deleted_files > 0:
        update_files_index_jsonl(done_embedding, unmodified, None, bucket, user_prefix, s3client)
    else:
        print(f"update_index_for_user_dropbox: No new embeddings or deleted files. Not updating files_index.jsonl")
    if time_limit_exceeded:
        print("update_index_for_user_dropbox: time limit exceeded. Forcing full listing next time..")
        dropbox_next_page_token = None # force full scan next time
    if partial_present(done_embedding):
        print("update_index_for_user_dropbox: partial present. Forcing full listing next time..")
        dropbox_next_page_token = None # force full scan next time
        
    # at this point, we must have the index_metadata.json file: it must have been downloaded above (and is unmodified as no files need to be embedded) or it was modified above (as files needed to be embedded)    
    _build_and_save_faiss_if_needed(email, None, s3client, bucket, prefix, sub_prefix, user_prefix, DocStorageType.DropBox )
    return dropbox_next_page_token
    
FOLDER = 'application/vnd.google-apps.folder'

def _process_gdrive_items(gdrive_listing, folder_details, items):
    if not items:
        print("No more files found.")
        return
    for item in items:
        filename = item['name']
        fileid = item['id']
        if 'size' in item:
            size = int(item['size'])
        else:
            size = 0
        mtime = from_rfc3339(item['modifiedTime'])
        mimetype = item['mimeType']
        if mimetype == 'application/vnd.google-apps.folder':
            if 'parents' in item:
                folder_details[fileid] = {'filename': filename, 'parents': item['parents']}
        else:
            if 'parents' in item:
                gdrive_listing[fileid] = {"filename": filename, "fileid": fileid, "mtime": mtime, 'mimetype': mimetype, 'size': size, 'parents': item['parents']}
            else:
                gdrive_listing[fileid] = {"filename": filename, "fileid": fileid, "mtime": mtime, 'mimetype': mimetype, 'size': size}
    return fileid

def _get_start_page_token(item):
    resp = None
    try:
        headers={"Authorization": f"Bearer {item['access_token']['S']}"}
        params={'supportsAllDrives': True}
        print(f"_get_start_page_token: startPageToken. hdrs={json.dumps(headers)}, params={json.dumps(params)}")
        resp = requests.get('https://www.googleapis.com/drive/v3/changes/startPageToken', headers=headers, params=params)
        resp.raise_for_status()
        respj = resp.json()
        print(f"_get_start_page_token: got {respj} from changes.getStartPageToken")
        return respj['startPageToken']
    except Exception as ex:
        if resp:
            print(f"_get_start_page_token: In changes.getStartPageToken for user {item['email']['S']}, caught {ex}. Response={resp.content}")
        else:
            print(f"_get_start_page_token: In changes.getStartPageToken for user {item['email']['S']}, caught {ex}")
        return None

def _get_gdrive_listing_incremental(service, s3_index, item, gdrive_next_page_token):
    status = False
    gdrive_listing = {}
    folder_details = {}
    deleted_files = {}

    new_start_page_token = None
    next_token = gdrive_next_page_token
    try:
        driveid = service.files().get(fileId='root').execute()['id']
        folder_details[driveid] = {'filename': 'My Drive'}
    except Exception as ex:
        print(f"_get_gdrive_listing_incremental: Exception {ex} while getting driveid")
        return status, gdrive_listing, deleted_files, new_start_page_token if new_start_page_token else next_token

    escape = 0
    while not new_start_page_token:
        escape += 1
        if escape > 10:
            break
        resp = None
        try:
            headers = {"Authorization": f"Bearer {item['access_token']['S']}", "Content-Type": "application/json"}
            params={'includeCorpusRemovals': True,
                'includeItemsFromAllDrives': True,
                'includeRemoved': True,
                'pageToken': next_token,
                'pageSize': 100,
                'restrictToMyDrive': False,
                'spaces': 'drive',
                'supportsAllDrives': True}
            print(f"_get_gdrive_listing_incremental: hdrs={json.dumps(headers)}, params={json.dumps(params)}")
            resp = requests.get('https://www.googleapis.com/drive/v3/changes', headers=headers, params=params)
            resp.raise_for_status()
            respj = resp.json()
            print(f"respj={respj}")
            for chg in respj['changes']:
                if chg['removed']:
                    if chg['changeType'] == 'file':
                        print(f"_get_gdrive_listing_incremental: removed. file. fileid={chg['fileId']}")
                        fileid = chg['fileId']
                        mtime = from_rfc3339(chg['time'])
                        deleted_files[fileid] = {"fileid": fileid, "mtime": mtime}
                else:
                    if chg['changeType'] == 'file':
                        filename = chg['file']['name']
                        fileid = chg['fileId']
                        mimetype = chg['file']['mimeType']
                        file_details = service.files().get(fileId=fileid, fields='size,parents,trashed,explicitlyTrashed,modifiedTime').execute()
                        print(f"_get_gdrive_listing_incremental: filename={filename}, fileid={fileid}, file_details={file_details}")
                        mtime = from_rfc3339(file_details['modifiedTime'])
                        if mimetype == 'application/vnd.google-apps.folder':
                            print(f"_get_gdrive_listing_incremental: changed. folder. filename={filename}, fileid={fileid}")
                            if 'parents' in file_details:
                                folder_details[fileid] = {'filename': filename, 'parents': file_details['parents']}
                        else:
                            size = file_details['size'] if 'size' in file_details else 0
                            if size: # ignore if size is not available
                                if file_details['trashed'] or file_details['explicitlyTrashed']:
                                    print(f"_get_gdrive_listing_incremental: filename={filename}, fileid={fileid}. trashed!")
                                    deleted_files[fileid] = {"fileid": fileid, "mtime": mtime}
                                else:
                                    if fileid in s3_index and mtime == s3_index[fileid]['mtime']:
                                        print(f"_get_gdrive_listing_incremental: filename={filename}, fileid={fileid}. unmodified since file details mtime matches s3index")
                                    else:
                                        if fileid in s3_index:
                                            print(f"_get_gdrive_listing_incremental: filename={filename}, fileid={fileid}. mtime={mtime}, s3_index_mtime={s3_index[fileid]['mtime']}")
                                        else:
                                            print(f"_get_gdrive_listing_incremental: filename={filename}, fileid={fileid}. mtime={mtime}, s3_index does not have fileid")
                                        entry = {"filename": filename, "fileid": fileid, "mtime": mtime,
                                                            'mimetype': mimetype, 'size': size}
                                        if 'parents' in file_details:
                                            entry['parents'] = file_details['parents']
                                        entry['path']=_calc_path(service, entry, folder_details)
                                        print(f"_get_gdrive_listing_incremental: changed. file. filename={filename}, fileid={fileid}, path={entry['path']}, mtime={mtime}")
                                        gdrive_listing[fileid] = entry
            if 'newStartPageToken' in respj:
                new_start_page_token = respj['newStartPageToken']
            else:
                next_token = respj['nextPageToken']
            status = True
        except Exception as ex:
            if resp:
                print(f"_get_gdrive_listing_incremental: caught {ex}")
                print(f"_get_gdrive_listing_incremental: status={resp.status_code}")
                print(f"_get_gdrive_listing_incremental: content={str(resp.content)}")
                print(f"_get_gdrive_listing_incremental: headers={str(resp.headers)}")
            else:
                print(f"_get_gdrive_listing_incremental: caught {ex}")
            status = False
    return status, gdrive_listing, deleted_files, new_start_page_token if new_start_page_token else next_token

def _get_gdrive_listing_full(service, item):
    gdrive_listing = {}
    pageToken = None
    kwargs = {}
    folder_details = {}
    gdrive_next_page_token = None
    try:
        driveid = service.files().get(fileId='root').execute()['id']
        folder_details[driveid] = {'filename': 'My Drive'}
    except Exception as ex:
        print(f"_get_gdrive_listing_full: Exception {ex} while getting driveid")
        return gdrive_listing, folder_details, gdrive_next_page_token
    
    total_items = 0
    while True:
        # Metadata for files: https://developers.google.com/drive/api/reference/rest/v3/files
        results = None
        for retry_num in range(1, 4):
            try:
                results = (
                    service.files()
                        .list(pageSize=100, orderBy="folder,modifiedTime desc, name",
                                fields="nextPageToken, files(id, name, size, modifiedTime, mimeType, parents)", **kwargs)
                    .execute()
                )
                break
            except Exception as ex:
                print(f"Caught {ex} while getting listing. retry_num={retry_num}")
                time.sleep(retry_num * 5)
                results=None
                continue
        if not results:
            print(f"_get_gdrive_listing_full: retried three times unsuccessfully. Returning without finishing..")
            break
        items = results.get("files", [])
        total_items += len(items)
        print(f"Fetched {total_items} items from google drive rooted at id={driveid}..")

        _process_gdrive_items(gdrive_listing, folder_details, items)

        pageToken = results.get("nextPageToken", None)
        kwargs['pageToken']=pageToken
        if not pageToken:
            gdrive_next_page_token = _get_start_page_token(item)
            print(f"_get_gdrive_listing_full: no pageToken. next_page_token={gdrive_next_page_token}")
            break
    return gdrive_listing, folder_details, gdrive_next_page_token

def _get_details_recursively(path, parent_fileid, file_details, dir_details):
    print(f"_get_detail_recursively: path={path}, parent_fileid={parent_fileid}")
    try:
        entries = os.listdir(path)
    except PermissionError as e:
        print(f"Permission denied: {path} - {e}")
        return
    except FileNotFoundError as e:
        print(f"Path not found: {path} - {e}")
        return
    for entry in entries:
        entry_path = os.path.join(path, entry)
        try:
            entry_stat = os.stat(entry_path)
            #mtime = DatetimeWithNanoseconds(datetime.datetime.fromtimestamp(entry_stat.st_mtime))
            mtime = from_microseconds(entry_stat.st_mtime*1000000.0)
            fileid = str(entry_stat.st_ino)
            if os.path.isdir(entry_path):
                dir_details[fileid] = {"filename": entry, "fileid": fileid,
                                        "mtime": mtime, "mimetype": 'application/vnd.google-apps.folder',
                                        "parents": [parent_fileid]}
                print(f"_get_detail_recursively: adding dir {entry_path}, fileid={fileid}")
                _get_details_recursively(entry_path, fileid, file_details, dir_details)
            else:
                file_details[fileid] = {"filename": entry, "fileid": fileid,
                                        "mtime": mtime, "mimetype": mimetypes.guess_type(entry)[0],
                                        "parents": [parent_fileid], "size": entry_stat.st_size}
                print(f"_get_detail_recursively: adding file {entry_path}, fileid={fileid}, size={entry_stat.st_size}")
        except PermissionError as e:
            print(f"Permission denied: {entry_path} - {e}")
        except FileNotFoundError as e:
            print(f"File not found: {entry_path} - {e}")
        except Exception as e:
            print(f"Error reading entry: {entry_path} - {e}")

def update_index_for_user_local(email, index_dir, docs_dir):
    print(f'update_index_for_user_local: Entered. email={email}, index_dir={index_dir}, docs_dir={docs_dir}')
    docs_dir_stat = os.stat(docs_dir)
    docs_dir_fileid = str(docs_dir_stat.st_ino)
    file_details = {}
    dir_details = {docs_dir_fileid: {"filename": docs_dir, "fileid": docs_dir_fileid}} # we use 1 as the dummy fileid for the docs_dir
    _get_details_recursively(docs_dir, docs_dir_fileid, file_details, dir_details)
    print(f"update_index_for_user_local: dir_details={dir_details}, file_details={file_details}")
    s3_index:Dict[str, dict] = get_s3_index(index_dir, None, None, None)
    print(f"update_index_for_user_local: s3_index={s3_index}")
    unmodified, needs_embedding, deleted_files = calc_file_lists(None, s3_index, file_details, dir_details)
    # remove deleted files from the index
    for fileid in deleted_files:
        if fileid in unmodified:
            print(f"update_index_for_user_local: removing deleted fileid {fileid} from s3_index")
            unmodified.pop(fileid)
        else:
            print(f"update_index_for_user_local: deleted fileid {fileid} not in unmodified index")

    # Move items in 'partial' state from umodified to needs_embedding
    to_del=[]
    for fid, entry in unmodified.items():
        if 'partial' in entry:
            pth = entry['path'] if 'path' in entry else ''
            print(f"update.gdrive: fileid={fid}, path={pth}, filename={entry['filename']} is in partial state. Moving from unmodified to needs_embedding")
            needs_embedding[fid] = entry
            to_del.append(fid)
    for td in to_del:
        del unmodified[td]

    gdrive_next_page_token = "notused_gdrive_next_token"
    done_embedding, time_limit_exceeded = process_files(email, LocalFileReader(docs_dir, file_details, dir_details),
                                                        unmodified, needs_embedding, index_dir, None, None, None)
    if len(done_embedding.items()) > 0 or len(deleted_files) > 0:
        update_files_index_jsonl(done_embedding, unmodified, index_dir, None, "", None)
    else:
        print(f"update_index_for_user_local: No new embeddings or deleted files. Not updating files_index.jsonl")
    if time_limit_exceeded:
        print("update_index_for_user_local: time limit exceeded. Forcing full listing next time..")
        gdrive_next_page_token = None # force full scan next time
    if partial_present(done_embedding):
        print("update_index_for_user_local: partial present. Forcing full listing next time..")
        gdrive_next_page_token = None # force full scan next time
    _build_and_save_faiss_if_needed(email, index_dir, None, None, None, None, "", DocStorageType.Local)
    return gdrive_next_page_token

def update_index_for_user_gdrive(email, s3client, bucket:str, prefix:str, gdrive_next_page_token:str=None):
    print(f'update_index_for_user_gdrive: Entered. {email}, gdrive_next_page_token={gdrive_next_page_token}')
    item = get_user_table_entry(email)
    # index1/xyz@abc.com
    user_prefix = f"{prefix}/{email}"
    try:
        # user_prefix = 'index1/raj@yoja.ai' 
        s3_index:Dict[str, dict] = get_s3_index(None, s3client, bucket, user_prefix)
        try:
            creds:google.oauth2.credentials.Credentials = refresh_user_google(item)
        except Exception as ex:
            print(f"update_index_for_user_gdrive: credentials not valid. not processing user {item['email']['S']}. Exception={ex}")
            return
        
        unmodified:dict; needs_embedding:dict;  deleted_files: dict
        service:googleapiclient.discovery.Resource = build("drive", "v3", credentials=creds)
        if not s3_index or not gdrive_next_page_token or gdrive_next_page_token == "1" or 'YOJA_FORCE_FULL_INDEX' in os.environ:
            if 'YOJA_FORCE_FULL_INDEX' in os.environ:
                print("update_index_for_user_gdrive: Forcing full index because of env var YOJA_FORCE_FULL_INDEX. Also setting gdrive_next_page_token to 1")
                gdrive_next_page_token = "1"
            gdrive_listing, folder_details, gdrive_next_page_token = _get_gdrive_listing_full(service, item)
            unmodified, needs_embedding, deleted_files = calc_file_lists(service, s3_index, gdrive_listing, folder_details)
            print(f"gdrive full_listing. Number of unmodified={len(unmodified)}; modified or added={len(needs_embedding)}; deleted={len(deleted_files)}; gdrive_next_page_token={gdrive_next_page_token}")
        else:
            status, needs_embedding, deleted_files, gdrive_next_page_token = _get_gdrive_listing_incremental(service, s3_index, item, gdrive_next_page_token)
            if status:
                unmodified = s3_index
            else:
                print("gdrive incremental listin failed. trying full_listing")
                gdrive_listing, folder_details, gdrive_next_page_token = _get_gdrive_listing_full(service, item)
                unmodified, needs_embedding, deleted_files = calc_file_lists(service, s3_index, gdrive_listing, folder_details)
                print(f"gdrive full_listing. Number of unmodified={len(unmodified)}; modified or added={len(needs_embedding)}; deleted={len(deleted_files)}; gdrive_next_page_token={gdrive_next_page_token}")
        # remove deleted files from the index
        for fileid in deleted_files:
            if fileid in unmodified:
                print(f"update_index_for_user_gdrive: removing deleted fileid {fileid} from s3_index")
                unmodified.pop(fileid)
            else:
                print(f"update_index_for_user_gdrive: deleted fileid {fileid} not in unmodified index")

        # Move items in 'partial' state from umodified to needs_embedding
        to_del=[]
        for fid, entry in unmodified.items():
            if 'partial' in entry:
                pth = entry['path'] if 'path' in entry else ''
                print(f"update.gdrive: fileid={fid}, path={pth}, filename={entry['filename']} is in partial state. Moving from unmodified to needs_embedding")
                needs_embedding[fid] = entry
                to_del.append(fid)
        for td in to_del:
            del unmodified[td]

        done_embedding, time_limit_exceeded = process_files(email, GoogleDriveReader(service, item['access_token']['S']),
                                                            unmodified, needs_embedding, None, s3client, bucket, user_prefix)
        if len(done_embedding.items()) > 0 or len(deleted_files) > 0:
            update_files_index_jsonl(done_embedding, unmodified, None, bucket, user_prefix, s3client)
        else:
            print(f"update_index_for_user_gdrive: No new embeddings or deleted files. Not updating files_index.jsonl")
        if time_limit_exceeded:
            print("update_index_for_user_gdrive: time limit exceeded. Forcing full listing next time..")
            gdrive_next_page_token = None # force full scan next time
        if partial_present(done_embedding):
            print("update_index_for_user_gdrive: partial present. Forcing full listing next time..")
            gdrive_next_page_token = None # force full scan next time

        # at this point, we must have the index_metadata.json file: it must have been downloaded above (and is unmodified as no files need to be embedded) or it was modified above (as files needed to be embedded)
        _build_and_save_faiss_if_needed(email, None, s3client, bucket, prefix, None, user_prefix, DocStorageType.GoogleDrive)
            
    except HttpError as error:
        print(f"HttpError occurred: {error}")
    except Exception as ex:
        print(f"Exception occurred: {ex}")
    return gdrive_next_page_token

def update_index_for_user(email, index_dir, docs_dir, s3client, bucket:str, prefix:str, start_time:datetime.datetime,
                            gdrive_next_page_token:str=None, dropbox_next_page_token:str=None):
    if index_dir:
        update_index_for_user_local(email, index_dir, docs_dir)
        return 'notused_gdrive_next_token', 'notused_dropbox_next_token'
    else:
        if not lambda_timelimit_exceeded():
            gdrive_next_page_token = update_index_for_user_gdrive(email, s3client, bucket, prefix, gdrive_next_page_token)
        if not lambda_timelimit_exceeded():
            dropbox_next_page_token = update_index_for_user_dropbox(email, s3client, bucket, prefix, dropbox_next_page_token)
        return gdrive_next_page_token, dropbox_next_page_token

if __name__=="__main__":
    with open(sys.argv[1], 'rb') as f:
        bio = io.BytesIO(f.read())
    vectorizer = get_vectorizer(email, [])
    rv = read_pdf(None, None, sys.argv[1], 'abc', bio, datetime.datetime.now(), vectorizer, [])
    rv['mtime'] = to_rfc3339(rv['mtime'])
    print(json.dumps(rv, indent=4))
    sys.exit(0)
